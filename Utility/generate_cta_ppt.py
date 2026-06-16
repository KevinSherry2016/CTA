from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


OUTPUT_PATH = Path(r"d:\CTA\CTA投研框架_9页_精修_v8.pptx")

SLIDE_W = 13.333
SLIDE_H = 7.5

NAVY = RGBColor(19, 43, 74)
BLUE = RGBColor(47, 95, 152)
TEAL = RGBColor(39, 118, 109)
GOLD = RGBColor(172, 118, 38)
RED = RGBColor(182, 79, 66)
INK = RGBColor(50, 57, 68)
MUTED = RGBColor(111, 120, 133)
BG = RGBColor(247, 248, 250)
CARD = RGBColor(255, 255, 255)
LINE = RGBColor(224, 229, 236)
LIGHT_BLUE = RGBColor(232, 239, 248)
LIGHT_TEAL = RGBColor(232, 243, 240)
LIGHT_GOLD = RGBColor(248, 240, 227)
LIGHT_RED = RGBColor(248, 235, 232)


def build_presentation() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(SLIDE_W)
    prs.slide_height = Inches(SLIDE_H)
    return prs


def add_background(slide, accent: RGBColor = NAVY) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG

    bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(SLIDE_W),
        Inches(0.12),
    )
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()


def add_title(slide, title: str, subtitle: str = "") -> None:
    title_box = slide.shapes.add_textbox(Inches(0.55), Inches(0.38), Inches(8.6), Inches(0.7))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(26)
    p.font.bold = True
    p.font.color.rgb = NAVY

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.58), Inches(1.02), Inches(11.5), Inches(0.38))
        tf = sub_box.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(12.5)
        p.font.color.rgb = MUTED


def add_footer(slide, page_no: int) -> None:
    box = slide.shapes.add_textbox(Inches(11.95), Inches(7.04), Inches(0.8), Inches(0.2))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = f"{page_no:02d}"
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(10)
    p.font.color.rgb = MUTED
    p.alignment = PP_ALIGN.RIGHT


def add_section_tag(slide, label: str, color: RGBColor = NAVY) -> None:
    tag = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(10.95), Inches(0.3), Inches(1.8), Inches(0.38)
    )
    tag.fill.solid()
    tag.fill.fore_color.rgb = color
    tag.line.fill.background()

    tf = tag.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = label
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = CARD
    p.alignment = PP_ALIGN.CENTER


def add_card(slide, x, y, w, h, title, bullets, accent=BLUE, fill_color=CARD, title_size=18, body_size=13, line_space=5):
    card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(h),
    )
    card.fill.solid()
    card.fill.fore_color.rgb = fill_color
    card.line.color.rgb = LINE
    card.line.width = Pt(1)

    top_bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(x),
        Inches(y),
        Inches(w),
        Inches(0.16),
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = accent
    top_bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(x + 0.22), Inches(y + 0.28), Inches(w - 0.44), Inches(0.42))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(title_size)
    p.font.bold = True
    p.font.color.rgb = accent

    body = slide.shapes.add_textbox(Inches(x + 0.22), Inches(y + 0.84), Inches(w - 0.4), Inches(h - 1.02))
    tf = body.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(2)
    tf.margin_right = Pt(2)
    tf.margin_top = Pt(1)
    tf.margin_bottom = Pt(1)
    tf.clear()
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = bullet
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(body_size)
        p.font.color.rgb = INK
        p.bullet = True
        p.space_after = Pt(line_space)


def add_labeled_card(slide, x, y, w, h, title, entries, accent=BLUE, fill_color=CARD, title_size=18, body_size=13, row_gap=0.16, label_w=0.95):
    card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h),
    )
    card.fill.solid()
    card.fill.fore_color.rgb = fill_color
    card.line.color.rgb = LINE
    card.line.width = Pt(1)

    top_bar = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(0.16),
    )
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = accent
    top_bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(x + 0.22), Inches(y + 0.28), Inches(w - 0.44), Inches(0.42))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(title_size)
    p.font.bold = True
    p.font.color.rgb = accent

    start_y = y + 0.86
    row_h = 0.68
    for idx, (label, value) in enumerate(entries):
        row_y = start_y + idx * row_h

        label_box = slide.shapes.add_textbox(Inches(x + 0.24), Inches(row_y), Inches(label_w), Inches(0.32))
        tf_label = label_box.text_frame
        tf_label.word_wrap = False
        tf_label.clear()
        p_label = tf_label.paragraphs[0]
        p_label.text = f"{label}："
        p_label.font.name = "Microsoft YaHei"
        p_label.font.size = Pt(body_size)
        p_label.font.bold = True
        p_label.font.color.rgb = INK

        value_box = slide.shapes.add_textbox(Inches(x + 0.20 + label_w), Inches(row_y), Inches(w - label_w - 0.42), Inches(0.52 + row_gap))
        tf_value = value_box.text_frame
        tf_value.word_wrap = True
        tf_value.margin_left = 0
        tf_value.margin_right = 0
        tf_value.margin_top = 0
        tf_value.margin_bottom = 0
        tf_value.clear()
        p_value = tf_value.paragraphs[0]
        p_value.text = value
        p_value.font.name = "Microsoft YaHei"
        p_value.font.size = Pt(body_size)
        p_value.font.color.rgb = INK


def add_callout(slide, x, y, w, h, text, fill_color=LIGHT_BLUE, font_color=NAVY):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    box = slide.shapes.add_textbox(Inches(x + 0.18), Inches(y + 0.1), Inches(w - 0.36), Inches(h - 0.2))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = font_color
    p.alignment = PP_ALIGN.CENTER


def style_line_series(series, color: RGBColor, width_pt: float = 2.5):
    series.format.line.color.rgb = color
    series.format.line.width = Pt(width_pt)


def slide_cover(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(250, 250, 248)

    left_panel = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(4.35), prs.slide_height
    )
    left_panel.fill.solid()
    left_panel.fill.fore_color.rgb = NAVY
    left_panel.line.fill.background()

    for idx, color in enumerate((BLUE, TEAL, GOLD)):
        band = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            Inches(0.55), Inches(1.2 + idx * 0.32), Inches(2.4), Inches(0.12)
        )
        band.fill.solid()
        band.fill.fore_color.rgb = color
        band.line.fill.background()

    t1 = slide.shapes.add_textbox(Inches(0.58), Inches(2.05), Inches(3.0), Inches(1.2))
    tf = t1.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "CTA\n投研框架"
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = CARD

    t2 = slide.shapes.add_textbox(Inches(4.9), Inches(1.5), Inches(6.8), Inches(0.65))
    tf = t2.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "从数据、因子到组合与执行"
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = NAVY

    body = slide.shapes.add_textbox(Inches(4.95), Inches(2.4), Inches(6.8), Inches(2.0))
    tf = body.text_frame
    tf.clear()
    lines = [
        "CTA 投研框架总览",
        "聚焦策略研发的关键链路：数据、因子、评估、组合、风控、执行",
        "目标：构建稳定、可扩展、可解释的商品期货量化体系",
    ]
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(15 if idx == 0 else 13)
        p.font.color.rgb = INK
        p.space_after = Pt(10)

    add_callout(slide, 4.95, 5.15, 2.25, 0.65, "策略主线")
    add_callout(slide, 7.45, 5.15, 2.25, 0.65, "组合增益")
    add_callout(slide, 9.95, 5.15, 2.25, 0.65, "风险约束")
    add_footer(slide, 1)


def slide_process(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_section_tag(slide, "流程", BLUE)
    add_title(slide, "CTA策略投研流程", "将研究链路拆分为五个环节，便于复盘与归因")

    steps = [
        ("数据", BLUE),
        ("因子", TEAL),
        ("投资组合", GOLD),
        ("风控", RED),
        ("执行", NAVY),
    ]
    x_positions = [0.65, 3.1, 5.55, 8.0, 10.45]
    for idx, ((name, color), x) in enumerate(zip(steps, x_positions)):
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(2.55), Inches(1.95), Inches(1.2)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = color
        shape.line.fill.background()
        tf = shape.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = name
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = CARD
        p.alignment = PP_ALIGN.CENTER
        if idx < len(steps) - 1:
            conn = slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                Inches(x + 1.95), Inches(3.15), Inches(x + 2.45), Inches(3.15)
            )
            conn.line.color.rgb = RGBColor(173, 179, 188)
            conn.line.width = Pt(2)

    add_card(
        slide, 0.75, 4.45, 12.0, 1.55, "流程要点",
        [
            "收益、风控与执行需要协同，核心壁垒在于全链路联动。",
            "任何单环节失真，最终都会表现为回撤放大、容量受限或执行偏差。",
        ],
        accent=NAVY,
        fill_color=CARD,
        title_size=16,
        body_size=12,
    )
    add_footer(slide, 2)


def slide_data(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_section_tag(slide, "数据", TEAL)
    add_title(slide, "数据体系", "数据覆盖决定信号来源，数据质量决定策略上限")
    add_card(slide, 0.6, 1.8, 4.05, 4.7, "量价数据", [
        "内容：K线、行情快照、Tick数据",
        "来源：交易所、Wind、Tushare",
        "用途：构建趋势、反转、波动率、量仓共振与微观结构类因子",
        "更新频率：覆盖毫秒级到日频",
    ], accent=BLUE, title_size=19, body_size=13, line_space=8)
    add_card(slide, 4.65, 1.8, 4.05, 4.7, "基本面数据", [
        "内容：库存、开工率、现货价格、CPI、利率、美元指数、信用利差",
        "来源：卓创资讯、钢联、上海有色网",
        "用途：支持品种强弱排序、行业判断以及宏观环境识别",
        "更新频率：日/周/月/季/年都覆盖",
    ], accent=TEAL, title_size=19, body_size=13, line_space=8)
    add_card(slide, 8.7, 1.8, 4.05, 4.7, "另类数据", [
        "内容：卫星图像、分析师预测、各类第三方事件数据",
        "来源：第三方数据供应商",
        "用途：补充传统数据难覆盖的数据，同时覆盖各类事件扰动与市场预期变化",
        "更新频率：日/周/月/季/年都覆盖",
    ], accent=GOLD, title_size=19, body_size=13, line_space=8)
    add_footer(slide, 3)


def slide_factors(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_section_tag(slide, "因子", GOLD)
    add_title(slide, "因子构建", "将复杂市场拆解为可量化表达的信号模块")
    add_card(slide, 0.65, 1.85, 3.95, 4.9, "量价因子", [
        "分类：趋势类、反转类、微结构、高频类",
        "示例：双均线、收益率偏度、能量潮、订单簿不平衡、微价格偏差",
        "优点：更新频率高、响应速度快、可直接映射交易行为",
        "适配性：适合趋势主导、波动放大、情绪切换较快的阶段",
        "特点：对行情拐点敏感，但对噪声和交易成本也更敏感",
    ], accent=BLUE, fill_color=LIGHT_BLUE, line_space=8)
    add_card(slide, 4.7, 1.85, 3.95, 4.9, "基本面因子", [
        "分类：产业基本面因子、宏观因子",
        "示例：库存变化、基差结构、CPI、PPI、利率、美元指数",
        "优点：经济含义清晰、跨周期解释力强、稳定性较好",
        "适配性：适合中低频方向判断和跨品种横向配置",
        "特点：更新频率相对较低，通常需要与量价信号联合使用",
    ], accent=TEAL, fill_color=LIGHT_TEAL, line_space=8)
    add_card(slide, 8.75, 1.85, 3.95, 4.9, "另类因子", [
        "分类：事件驱动类、替代数据类、预期偏差类",
        "示例：天气扰动、卫星产能监测、产业链事件映射",
        "优点：信息增量高，能捕捉传统因子难覆盖的变化",
        "适配性：适合突发事件驱动或常规信号拥挤阶段",
        "特点：噪声较高，需先做清洗和验证，再进入生产流程",
    ], accent=GOLD, fill_color=LIGHT_GOLD, line_space=8)
    add_footer(slide, 4)


def slide_evaluation(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_section_tag(slide, "评估", RED)
    add_title(slide, "评估体系", "不仅评估单因子收益，还评估其对现有组合的边际增益")

    add_card(slide, 0.65, 1.75, 2.65, 4.9, "基础指标", [
        "夏普比率",
        "最大回撤深度",
        "最大回撤时间",
        "持仓周期",
        "容量约束",
        "滑点敏感度",
    ], accent=BLUE)

    # 约3年 月度时间轴
    dates_long = [
        "2024-01", "2024-02", "2024-03", "2024-04", "2024-05", "2024-06",
        "2024-07", "2024-08", "2024-09", "2024-10", "2024-11", "2024-12",
        "2025-01", "2025-02", "2025-03", "2025-04", "2025-05", "2025-06",
        "2025-07", "2025-08", "2025-09", "2025-10", "2025-11", "2025-12",
        "2026-01", "2026-02", "2026-03", "2026-04", "2026-05", "2026-06",
    ]
    # 仅保留每3个月一个横轴标签，避免时间坐标过密。
    dates_sparse = [d if i % 3 == 0 else "" for i, d in enumerate(dates_long)]
    gmv_vals = [2.8, 3.2, 2.5, 3.6, 2.3, 3.8, 3.1, 2.6, 3.4, 2.9, 3.7, 2.4,
                3.3, 2.7, 3.9, 2.2, 3.5, 2.8, 3.2, 2.5, 3.7, 3.0, 2.4, 3.6,
                2.9, 3.3, 2.6, 3.8, 2.3, 3.1]
    delta_vals = [-0.5, -0.3, 0.4, 0.1, 0.8, 1.2, 1.5, 0.9, 0.3, -0.4, -0.9, 0.2,
                  -0.7, -0.2, 0.6, 0.3, 1.0, 1.8, 2.1, 1.3, 0.6, -0.1, -0.8, 0.3,
                  -1.4, -0.8, 0.2, -0.3, 1.1, 2.0]

    # --- 左上：GMV时间序列 ---
    d_gmv = CategoryChartData()
    d_gmv.categories = dates_sparse
    d_gmv.add_series("GMV", gmv_vals)
    c_gmv = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE, Inches(3.55), Inches(1.8), Inches(4.85), Inches(2.25), d_gmv,
    ).chart
    c_gmv.has_title = True
    c_gmv.chart_title.text_frame.text = "GMV"
    c_gmv.has_legend = False
    c_gmv.value_axis.has_major_gridlines = False
    c_gmv.value_axis.minimum_scale = 0.0
    c_gmv.value_axis.maximum_scale = 5.0
    c_gmv.category_axis.tick_labels.font.size = Pt(7)
    c_gmv.value_axis.tick_labels.font.size = Pt(8)
    style_line_series(c_gmv.series[0], BLUE, 2.5)

    # --- 右上：Delta时间序列 ---
    d_delta = CategoryChartData()
    d_delta.categories = dates_sparse
    d_delta.add_series("Delta", delta_vals)
    c_delta = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE, Inches(8.45), Inches(1.8), Inches(4.35), Inches(2.25), d_delta,
    ).chart
    c_delta.has_title = True
    c_delta.chart_title.text_frame.text = "Delta"
    c_delta.has_legend = False
    c_delta.value_axis.has_major_gridlines = False
    c_delta.value_axis.minimum_scale = -2.0
    c_delta.value_axis.maximum_scale = 3.0
    c_delta.category_axis.tick_labels.font.size = Pt(7)
    c_delta.value_axis.tick_labels.font.size = Pt(8)
    style_line_series(c_delta.series[0], BLUE, 2.5)

    # --- 左下：PnL工作日分布（正负双色）---
    d_wd = CategoryChartData()
    d_wd.categories = ["周一", "周二", "周三", "周四", "周五"]
    wd_vals = [0.12, 0.08, -0.03, 0.15, 0.09]
    d_wd.add_series("PnL", wd_vals)
    c_wd = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(3.55), Inches(4.15), Inches(4.85), Inches(2.25), d_wd,
    ).chart
    c_wd.has_title = True
    c_wd.chart_title.text_frame.text = "PnL 工作日分布"
    c_wd.has_legend = False
    c_wd.value_axis.has_major_gridlines = False
    c_wd.value_axis.minimum_scale = -0.15
    c_wd.value_axis.maximum_scale = 0.25
    c_wd.category_axis.tick_labels.font.size = Pt(10)
    c_wd.value_axis.tick_labels.font.size = Pt(8)
    from pptx.oxml.ns import qn as _qn
    c_wd.category_axis._element.find(_qn('c:tickLblPos')).set('val', 'low')
    c_wd.series[0].format.fill.solid()
    c_wd.series[0].format.fill.fore_color.rgb = BLUE
    # 禁止负值反色，保持同一填充色
    from pptx.oxml.ns import qn
    from lxml import etree
    ser_el = c_wd.series[0]._element
    inv = ser_el.find(qn('c:invertIfNegative'))
    if inv is None:
        inv = etree.SubElement(ser_el, qn('c:invertIfNegative'))
    inv.set('val', '0')

    # --- 右下：PnL月份分布（正负双色）---
    d_mo = CategoryChartData()
    d_mo.categories = ["1月", "2月", "3月", "4月", "5月", "6月",
                        "7月", "8月", "9月", "10月", "11月", "12月"]
    mo_vals = [0.18, -0.05, 0.22, 0.07, -0.11, 0.14, 0.09, 0.03, 0.16, -0.08, 0.20, 0.11]
    d_mo.add_series("PnL", mo_vals)
    c_mo = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(8.45), Inches(4.15), Inches(4.35), Inches(2.25), d_mo,
    ).chart
    c_mo.has_title = True
    c_mo.chart_title.text_frame.text = "PnL 月份分布"
    c_mo.has_legend = False
    c_mo.value_axis.has_major_gridlines = False
    c_mo.value_axis.minimum_scale = -0.15
    c_mo.value_axis.maximum_scale = 0.30
    c_mo.category_axis.tick_labels.font.size = Pt(8)
    c_mo.value_axis.tick_labels.font.size = Pt(8)
    c_mo.category_axis._element.find(qn('c:tickLblPos')).set('val', 'low')
    c_mo.series[0].format.fill.solid()
    c_mo.series[0].format.fill.fore_color.rgb = BLUE
    ser_el_mo = c_mo.series[0]._element
    inv_mo = ser_el_mo.find(qn('c:invertIfNegative'))
    if inv_mo is None:
        inv_mo = etree.SubElement(ser_el_mo, qn('c:invertIfNegative'))
    inv_mo.set('val', '0')

    add_footer(slide, 5)


def slide_incremental_evaluation(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_section_tag(slide, "评估", RED)
    add_title(slide, "增量评估", "通过新因子累计与残差累计判断新增信号对原组合的边际贡献")

    chart_data = CategoryChartData()
    chart_data.categories = [
        "2021Q1", "2021Q2", "2021Q3", "2021Q4",
        "2022Q1", "2022Q2", "2022Q3", "2022Q4",
        "2023Q1", "2023Q2", "2023Q3", "2023Q4",
        "2024Q1", "2024Q2", "2024Q3", "2024Q4",
        "2025Q1", "2025Q2", "2025Q3", "2025Q4",
    ]
    chart_data.add_series("累计收益率", [
        0, 13, 28, 20,
        40, 58, 78, 68,
        90, 110, 133, 123,
        153, 175, 198, 183,
        215, 238, 260, 280,
    ])
    chart_data.add_series("累计残差", [
        0, 3, 6, 5,
        10, 14, 19, 17,
        22, 27, 33, 30,
        37, 43, 49, 46,
        54, 61, 68, 80,
    ])
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE,
        Inches(2.55), Inches(1.6), Inches(8.25), Inches(2.75),
        chart_data,
    ).chart
    chart.has_title = True
    chart.chart_title.text_frame.text = "累计收益率 / 累计残差"
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.value_axis.has_major_gridlines = False
    chart.value_axis.maximum_scale = 300.0
    chart.category_axis.tick_labels.number_format = "General"
    chart.category_axis.tick_labels.font.size = Pt(9)
    chart.value_axis.tick_labels.font.size = Pt(9)
    style_line_series(chart.series[0], BLUE, 3.0)
    style_line_series(chart.series[1], RED, 3.0)

    # 盘中（日内）平均走势：用 matplotlib 画真正垂直阶跃曲线，插入为图片
    import matplotlib
    matplotlib.use('Agg')
    import matplotlib.pyplot as plt
    import matplotlib.patches as mpatches
    import io

    times = [
        "21:00", "21:30", "22:00", "22:30", "23:00", "23:30",
        "09:00", "09:30", "10:00", "10:30", "11:00", "11:30",
        "13:30", "14:00", "14:30", "15:00",
    ]
    # 基础平稳值；21:00索引0和09:00索引6处做垂直跳跃
    base = [0.02, 0.28, 0.30, 0.29, 0.32, 0.34,
            0.35, 0.62, 0.61, 0.64, 0.63, 0.67,
            0.69, 0.72, 0.71, 0.75]

    fig, ax = plt.subplots(figsize=(9.5, 1.5))
    xs = list(range(len(times)))

    # 画平稳段：逐段绘制，在跳跃点画垂直线
    jump_indices = {0: (0.02, 0.25), 6: (0.34, 0.58)}
    prev_y = base[0]
    for i in range(len(times)):
        if i in jump_indices:
            lo, hi = jump_indices[i]
            # 垂直线
            ax.plot([i, i], [lo, hi], color='#AC7626', linewidth=2.0)
            prev_y = hi
        if i < len(times) - 1:
            next_y = base[i + 1] if (i + 1) not in jump_indices else jump_indices[i + 1][0]
            ax.plot([i, i + 1], [prev_y, next_y], color='#AC7626', linewidth=2.0)
            prev_y = next_y

    ax.set_xticks(xs)
    ax.set_xticklabels(times, fontsize=7, rotation=30, ha='right')
    ax.set_ylabel("", fontsize=8)
    ax.set_title("Intraday Average PnL", fontsize=9, pad=4)
    ax.set_ylim(-0.05, 1.35)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.tick_params(axis='y', labelsize=7)
    ax.yaxis.grid(False)
    fig.tight_layout(pad=0.3)

    buf = io.BytesIO()
    fig.savefig(buf, format='png', dpi=150)
    plt.close(fig)
    buf.seek(0)

    slide.shapes.add_picture(buf, Inches(2.55), Inches(4.55), Inches(8.25), Inches(1.95))

    add_footer(slide, 6)


def slide_portfolio(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_section_tag(slide, "组合", NAVY)
    add_title(slide, "投资组合", "组合不是简单叠加高Sharpe因子，而是在相关性约束下实现稳定输出")

    add_card(slide, 0.7, 1.8, 4.15, 4.9, "核心原则", [
        "波动率标准化后再合成权重",
        "新因子与现有组合须保持低相关",
        "评估标准是整体组合增益，而非单腿收益",
        "单策略权重设上限，防止局部主导",
    ], accent=BLUE, body_size=13)

    add_card(slide, 4.95, 1.8, 3.75, 4.9, "组合结构", [
        "趋势类 — 捕捉主线行情",
        "反转类 — 补足震荡阶段收益",
        "另类类 — 提供增量信息",
        "按Sector与品种分层配置",
        "按月复核风格暴露",
    ], accent=TEAL, body_size=13)
    add_card(slide, 8.85, 1.8, 3.75, 4.9, "落地判断", [
        "优先参数高原，避免参数尖峰",
        "关注净值平滑度与回撤修复速度",
        "相关性、容量与换手联合约束",
        "同等收益下优选低交易成本方案",
    ], accent=GOLD, body_size=13)
    add_footer(slide, 7)


def slide_risk(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_section_tag(slide, "风控", RED)
    add_title(slide, "风控", "风控目标不是压制收益，而是确保组合在极端环境下仍可管理")

    add_card(slide, 0.7, 1.8, 3.9, 2.1, "硬约束", [
        "VaR约束",
        "投资限制",
        "单品种与单板块集中度限制",
    ], accent=RED)
    add_card(slide, 0.7, 4.1, 3.9, 2.1, "异常情况处理", [
        "盘中剧烈波动时，需考虑无法按常规方式平仓的场景",
        "预设降频、减仓、替代合约与人工接管机制",
        "异常行情结束后，执行复盘并更新熔断参数",
    ], accent=GOLD)
    add_card(slide, 4.9, 1.8, 7.4, 4.5, "仓位系数管理", [
        "第一层：节假日前后，按流动性与跳空风险调整仓位",
        "第二层：触发预警线、止损线后，自动调整单策略仓位",
        "第三层：当某类策略持续回撤时，逐步下调该策略族权重",
    ], accent=NAVY)

    add_footer(slide, 8)


def slide_execution(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_section_tag(slide, "执行", BLUE)
    add_title(slide, "交易执行", "执行不是末端细节，而是影响策略收益兑现的关键环节")

    add_card(slide, 0.9, 1.9, 5.55, 4.8, "VWAP为主，TWAP为辅", [
        "适用场景：日度小幅调仓",
        "适用场景：常规权重调整",
        "核心目标：降低冲击成本",
        "核心目标：兼顾成交效率",
    ], accent=BLUE, body_size=14)
    add_card(slide, 6.85, 1.9, 5.55, 4.8, "TWAP为主，VWAP为辅", [
        "适用场景：建仓、清仓、换月",
        "适用场景：集中执行场景",
        "核心目标：平滑市场冲击",
        "核心目标：提高执行可控性",
    ], accent=TEAL, body_size=14)
    add_footer(slide, 9)


def slide_summary(prs: Presentation):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide)
    add_section_tag(slide, "总结", NAVY)
    add_title(slide, "总结", "以统一框架打通数据、研究与交易")

    add_card(slide, 0.75, 1.85, 12.0, 4.9, "关键结论", [
        "数据决定信号上限",
        "因子决定收益来源",
        "组合决定收益稳定性",
        "风控与执行决定收益能否兑现",
        "长期竞争力来自流程化、标准化与可复制性",
    ], accent=NAVY, title_size=20, body_size=15)
    add_footer(slide, 10)


if __name__ == "__main__":
    prs = build_presentation()
    slide_cover(prs)
    slide_process(prs)
    slide_data(prs)
    slide_factors(prs)
    slide_evaluation(prs)
    slide_incremental_evaluation(prs)
    slide_portfolio(prs)
    slide_risk(prs)
    slide_execution(prs)
    slide_summary(prs)
    prs.save(OUTPUT_PATH)
    print(f"Saved: {OUTPUT_PATH}")